import collections 
import collections.abc
from pptx import Presentation
import google.generativeai as palm
import config
import numpy as np

palm.configure(api_key=config.api_key)
embedding_model = "models/embedding-gecko-001"
text_model = "models/text-bison-001"

def tup_embeddings(tup):
    return (tup[1], (tup[0], palm.generate_embeddings(model=embedding_model, text=tup[1])['embedding']))

def test_powerpoint():
    prs = Presentation("testdocs/test.pptx")
    text_blobs = []
    for slide in prs.slides:
        text_blob = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_blob += shape.text.replace("\n", "\t")
        text_blobs.append(text_blob)
    text_blobs = [text_blob for text_blob in text_blobs if len(text_blob) > 10]
    indexed_text_blobs = enumerate(text_blobs)

    embedding_dict = dict(map(tup_embeddings, indexed_text_blobs))   
    
    while True:
        # Prompt, Generate, Score
        user_input = input("Enter a question: ")
        user_input_embedding = palm.generate_embeddings(model=embedding_model, text=user_input)['embedding']
        def score(embedding):
            return np.dot(user_input_embedding, embedding)
        scores = {text_blob: (place, score(embedding)) for text_blob, (place, embedding) in embedding_dict.items()}
        scores = sorted(scores.items(), key=lambda x: x[1][1], reverse=True)
        
        # Top 4 matches for background
        top_scores_inds = [ind for _, (ind, _) in scores[:4]]
        background = [text_blobs[i] for i in top_scores_inds]
        background = "\n".join(background)
        request = "\nUsing the above information, answer the following question, \"" + user_input + "\":\n"
        prompt = background + request
        print(prompt)
        
        # Generate response
        response = palm.generate_text(prompt=prompt, model=text_model)
        if response.result != None:
            print("Answer: ", response.result)
        else:
            print("No answer found, expanding search.")
            # Locality 
            top_score_indices = [index for _, (index, _) in scores[:2]]
            ranges = []
            for (i, index) in enumerate(top_score_indices):
                if i == 0:
                    start = index -1 if index - 1 >= 0 else 0
                    end = index + 1 if index + 1 < len(text_blobs) else len(text_blobs) - 1
                else:
                    if index - 1 <= end:
                        end = index + 1 if index + 1 < len(text_blobs) else len(text_blobs) - 1
                    else:
                        ranges.append((start, end + 1))
                        start = index - 1 if index - 1 >= 0 else 0
                        end = index + 1 if index + 1 < len(text_blobs) else len(text_blobs) - 1
            ranges.append((start, end + 1))
            
            background_list = []
            for (start, end) in ranges:
                background_list = [text_blob for text_blob in text_blobs[start:end]]
            
            background = "\n".join(background_list)
            prompt = background + request
            
            response = palm.generate_text(prompt=prompt, model=text_model)
            if response.result != None:
                print("Answer: ", response.result)
            else:
                print("No answer found.")
    

def test_pdf():
    pass
    
if __name__ == "__main__":
    test_powerpoint()